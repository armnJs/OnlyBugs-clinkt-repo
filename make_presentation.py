import sys
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# 16:9 widescreen dimensions
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color Palette (Dark Mode Premium)
BG_COLOR = RGBColor(15, 17, 23)        # #0F1117
CARD_BG = RGBColor(26, 29, 39)         # #1A1D27
CARD_BORDER = RGBColor(46, 49, 64)     # #2E3140
ACCENT = RGBColor(108, 99, 255)        # #6C63FF
GREEN = RGBColor(52, 211, 153)         # #34D399
RED = RGBColor(255, 90, 90)            # #FF5A5A
AMBER = RGBColor(251, 191, 36)         # #FBBF24
TEXT_MAIN = RGBColor(232, 233, 237)    # #E8E9ED
TEXT_MUTED = RGBColor(139, 143, 163)   # #8B8FA3

def set_slide_background(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BG_COLOR

def add_header(slide, title_text, subtitle_text, category_tag="CLINKT CASE INVESTIGATION"):
    tag_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(8), Inches(0.35))
    tf_tag = tag_box.text_frame
    tf_tag.word_wrap = True
    p_tag = tf_tag.paragraphs[0]
    p_tag.text = category_tag.upper()
    p_tag.font.size = Pt(11)
    p_tag.font.bold = True
    p_tag.font.color.rgb = ACCENT

    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.75), Inches(11.5), Inches(0.6))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = TEXT_MAIN

    if subtitle_text:
        sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.35), Inches(11.5), Inches(0.4))
        tf_sub = sub_box.text_frame
        tf_sub.word_wrap = True
        p_sub = tf_sub.paragraphs[0]
        p_sub.text = subtitle_text
        p_sub.font.size = Pt(13)
        p_sub.font.color.rgb = TEXT_MUTED

def add_card(slide, left, top, width, height, title, value=None, value_color=None, body_lines=None, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = CARD_BG
    shape.line.color.rgb = border_color if border_color else CARD_BORDER
    shape.line.width = Pt(1.5)
    
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_right = Inches(0.2)
    tf.margin_top = Inches(0.2)
    tf.margin_bottom = Inches(0.2)
    
    p0 = tf.paragraphs[0]
    p0.text = title.upper()
    p0.font.size = Pt(11)
    p0.font.bold = True
    p0.font.color.rgb = TEXT_MUTED
    
    if value:
        p_val = tf.add_paragraph()
        p_val.text = str(value)
        p_val.font.size = Pt(28)
        p_val.font.bold = True
        p_val.font.color.rgb = value_color if value_color else TEXT_MAIN
        p_val.space_after = Pt(6)
        
    if body_lines:
        for line in body_lines:
            p_line = tf.add_paragraph()
            p_line.text = line
            p_line.font.size = Pt(11.5)
            p_line.font.color.rgb = TEXT_MAIN
            p_line.space_after = Pt(3)

blank_layout = prs.slide_layouts[6]

# SLIDE 1: Title Slide
slide1 = prs.slides.add_slide(blank_layout)
set_slide_background(slide1)

tag_box = slide1.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(6), Inches(0.4))
p = tag_box.text_frame.paragraphs[0]
p.text = "EXECUTIVE CASE REPORT & LIVE SOLUTION"
p.font.size = Pt(13)
p.font.bold = True
p.font.color.rgb = ACCENT

title_box = slide1.shapes.add_textbox(Inches(1.0), Inches(2.3), Inches(11.3), Inches(1.8))
tf = title_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Where Clinkt's Revenue Goes to Die"
p.font.size = Pt(44)
p.font.bold = True
p.font.color.rgb = TEXT_MAIN

p2 = tf.add_paragraph()
p2.text = "Root-Cause Investigation of the 70.5% Cart Abandonment & Our Working +87.6% Revenue Recovery Engine"
p2.font.size = Pt(18)
p2.font.color.rgb = TEXT_MUTED
p2.space_before = Pt(12)

add_card(slide1, Inches(1.0), Inches(4.8), Inches(2.6), Inches(1.6), "Lost Cart Revenue", "₹84,857", RED, ["1.75× Actual Sales"], RED)
add_card(slide1, Inches(3.9), Inches(4.8), Inches(2.6), Inches(1.6), "Cart Abandonment", "70.5%", AMBER, ["666 of 945 dropped"], AMBER)
add_card(slide1, Inches(6.8), Inches(4.8), Inches(2.6), Inches(1.6), "Stock Issues", "60.7%", RED, ["Sub-optimal stock days"], RED)
add_card(slide1, Inches(9.7), Inches(4.8), Inches(2.6), Inches(1.6), "Projected Recovery", "+87.6%", GREEN, ["+₹42.4K/mo uplift"], GREEN)

# SLIDE 2: Core Verdict
slide2 = prs.slides.add_slide(blank_layout)
set_slide_background(slide2)
add_header(slide2, "The Executive Verdict: A Single Self-Feeding Failure Loop", "Traffic is arriving (2,708 views), but revenue leaks at every step of the funnel.")

add_card(slide2, Inches(0.8), Inches(2.0), Inches(3.6), Inches(4.8), "1. Recommendation Mismatch", "15.5% vs 9%", RED, [
    "• Home Care pushed heavily (15.5% views) but accounts for only 9.0% orders.",
    "• High-intent Staples & Grains buried (8.3% views vs 11.5% orders).",
    "• 264 sessions (32.6%) bounce in <1.4 min with zero cart additions."
])

add_card(slide2, Inches(4.8), Inches(2.0), Inches(3.6), Inches(4.8), "2. Blind Inventory Planning", "0 Days", RED, [
    "• 60.7% of product-days face stock shortages.",
    "• Top demand SKU (Instant Noodles, score 481) had 0 healthy days & 11 critical stockout days.",
    "• Interested customers add items, hit stockouts, and drop."
])

add_card(slide2, Inches(8.8), Inches(2.0), Inches(3.6), Inches(4.8), "3. High-Value Abandonment", "79.5%", AMBER, [
    "• Family segment has highest AOV (₹192.60) but worst cart drop rate (79.5%).",
    "• ₹84,857 left behind in carts vs ₹48,412 earned revenue.",
    "• Total order velocity trending down (r = -0.212)."
])

# SLIDE 3: Funnel Analysis
slide3 = prs.slides.add_slide(blank_layout)
set_slide_background(slide3)
add_header(slide3, "Funnel Leakage: 70.5% of Carts Abandoned", "Where customer intent is generated vs where it evaporates.")

add_card(slide3, Inches(0.8), Inches(2.0), Inches(5.6), Inches(4.8), "The 3-Step Funnel Breakdown", None, None, [
    "Step 1: Product Views (811 Sessions)",
    "  • 2,708 Views generated across 31 days (100% baseline)",
    "  • 65.1% drop before adding anything to cart",
    "",
    "Step 2: Add-to-Cart (547 Sessions)",
    "  • 945 Total Items added to cart (34.9% view conversion)",
    "  • 70.5% OF ALL CART ITEMS ARE ABANDONED (666 items)",
    "",
    "Step 3: Completed Orders (218 Sessions)",
    "  • 279 Total Items purchased (Only 29.5% cart capture)",
    "  • Actual Revenue: ₹48,412  |  Abandoned: ₹84,857"
])

add_card(slide3, Inches(6.8), Inches(2.0), Inches(5.7), Inches(2.25), "Family Segment Crisis", "79.5%", RED, [
    "• Family AOV: ₹192.60 (Highest on platform)",
    "• 174 abandoned items vs 45 completed orders",
    "• Losing our most valuable customer cohort."
], RED)

add_card(slide3, Inches(6.8), Inches(4.55), Inches(5.7), Inches(2.25), "Bounce Duration Disparity", "1.4m vs 9.4m", AMBER, [
    "• View-only sessions: Average 1.4 minutes (Instant bounce)",
    "• Completed order sessions: Average 9.4 minutes (6.7× longer)",
    "• Proof that discovery friction triggers early abandonment."
], AMBER)

# SLIDE 4: Scientific Testing & Ruled Out
slide4 = prs.slides.add_slide(blank_layout)
set_slide_background(slide4)
add_header(slide4, "Scientific Rigor: What We Tested & Ruled Out", "Every alternative explanation was systematically tested against dataset evidence.")

add_card(slide4, Inches(0.8), Inches(2.0), Inches(3.6), Inches(4.8), "1. Price Sensitivity?", "RULED OUT", GREEN, [
    "Hypothesis: Customers abandon because prices are too high.",
    "",
    "Evidence:",
    "• Abandoned items avg: ₹127.41",
    "• Ordered items avg: ₹123.63",
    "• Difference is only 3.0% (negligible).",
    "• Price is NOT driving the 70.5% drop."
], GREEN)

add_card(slide4, Inches(4.8), Inches(2.0), Inches(3.6), Inches(4.8), "2. Total Supply Shortage?", "RULED OUT", GREEN, [
    "Hypothesis: The warehouse had zero stock across the board.",
    "",
    "Evidence:",
    "• Slow movers (Cookies 200g, Cola 750ml) maintained healthy stock for 31/31 days.",
    "• This is an allocation/prioritization failure, not a supply ceiling."
], GREEN)

add_card(slide4, Inches(8.8), Inches(2.0), Inches(3.6), Inches(4.8), "3. Uninterested Visitors?", "RULED OUT", GREEN, [
    "Hypothesis: Website traffic is accidental or low-intent clicks.",
    "",
    "Evidence:",
    "• 100% of purchased categories were previously viewed in session.",
    "• When relevant items are shown and in stock, conversion occurs."
], GREEN)

# SLIDE 5: The 3-Pillar Solution Engine
slide5 = prs.slides.add_slide(blank_layout)
set_slide_background(slide5)
add_header(slide5, "The Solution Engine: 3 Working Algorithmic Models", "We engineered and exported actionable optimization models for all 36 SKUs.")

add_card(slide5, Inches(0.8), Inches(2.0), Inches(3.6), Inches(4.8), "Pillar 1: Dynamic Replenishment", "ROP Model", ACCENT, [
    "• Formula: ROP = (Latent_Demand × L) + Safety_Stock",
    "• Factors in uncaptured cart demand.",
    "• Instant Noodles buffer increased from 10 to 24 units.",
    "• Bath Soap buffer to 18 units.",
    "• Full 36-SKU plan in exports/13."
], ACCENT)

add_card(slide5, Inches(4.8), Inches(2.0), Inches(3.6), Inches(4.8), "Pillar 2: Recommendation Matrix", "Re-Weighted", ACCENT, [
    "• Suppress Home Care: -42% (down to 9.0%)",
    "• Suppress Personal Care: -20% (down to 9.3%)",
    "• Boost Staples & Grains: +66% (to 13.8%)",
    "• Boost Fresh Produce: +42% (to 15.5%)",
    "• Eliminates wasted landing impressions."
], ACCENT)

add_card(slide5, Inches(8.8), Inches(2.0), Inches(3.6), Inches(4.8), "Pillar 3: Family Basket Engine", "Retargeting", ACCENT, [
    "• Dedicated 'Family Essentials' smart packs.",
    "• Real-time stockout warnings before cart addition.",
    "• Automated SMS/Push cart recovery targeting 329 dropped sessions.",
    "• Recovers 20-35% of lost carts."
], ACCENT)

# SLIDE 6: Financial ROI & Live Deliverables
slide6 = prs.slides.add_slide(blank_layout)
set_slide_background(slide6)
add_header(slide6, "Financial Impact & Live Prototype Suite", "Quantified business recovery models and live interactive tools.")

add_card(slide6, Inches(0.8), Inches(2.0), Inches(2.6), Inches(2.2), "Current Revenue", "₹48,412", TEXT_MAIN, ["August Baseline"])
add_card(slide6, Inches(3.8), Inches(2.0), Inches(2.6), Inches(2.2), "Recovered Revenue", "+₹32,585", GREEN, ["From ₹84.9K cart pool"], GREEN)
add_card(slide6, Inches(6.8), Inches(2.0), Inches(2.6), Inches(2.2), "Projected Top-Line", "₹80,997", ACCENT, ["+67.3% Growth"], ACCENT)
add_card(slide6, Inches(9.8), Inches(2.0), Inches(2.6), Inches(2.2), "Max Potential", "₹90,841", GREEN, ["+87.6% (50% recovery)"], GREEN)

add_card(slide6, Inches(0.8), Inches(4.5), Inches(11.6), Inches(2.3), "Deliverables Ready for Live Demo", None, None, [
    "1. Interactive BI Dashboard & Simulator (dashboard.html): Live real-time sliders for recovery forecasting & 10 visual charts.",
    "2. Unified Storefront & Smart Cart SPA (clinkt-ui.html): Full product catalog, instant cart management, and Dark/Light mode.",
    "3. Executive Findings Report PDF (findings_report.pdf): Formal 7-page presentation document.",
    "4. 14 Cleaned & Optimized Datasets (cleaned data/exports/): Including replenishment schedules & recommendation weights."
])

output_path = r"c:\Users\VIBEHACK\Desktop\clinkt-repo\presentation.pptx"
prs.save(output_path)
print(f"Presentation saved to {output_path}")
